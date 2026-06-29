#!/usr/bin/env python3
"""
🧠 DereInside Deck — 40+ slides · Local-First AI Knowledge System
Safety Orange theme, matching the Derekcoding talk style.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

ORANGE = RGBColor(0xFF, 0x6B, 0x35)
DARK   = RGBColor(0x0A, 0x0A, 0x0A)
LIGHT  = RGBColor(0xFA, 0xFA, 0xF8)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
INK    = RGBColor(0x0A, 0x0A, 0x0A)
G1     = RGBColor(0xF0, 0xF0, 0xEE)
G2     = RGBColor(0xD4, 0xD4, 0xD2)
G3     = RGBColor(0x73, 0x73, 0x73)
BLUE   = RGBColor(0x38, 0xBD, 0xF8)
GREEN  = RGBColor(0x74, 0xE8, 0x9F)
RED    = RGBColor(0xF8, 0x71, 0x71)
PURPLE = RGBColor(0xA7, 0x8B, 0xFA)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
N = 0

def bg(s, c):
    f = s.background.fill; f.solid(); f.fore_color.rgb = c

def tb(s, l, t, w, h, txt, sz=18, c=INK, b=False, fn="Inter", al=PP_ALIGN.LEFT):
    bx = s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h))
    p = bx.text_frame.paragraphs[0]; p.text = txt
    p.font.size = Pt(sz); p.font.color.rgb = c; p.font.bold = b; p.font.name = fn; p.alignment = al
    return bx

def mtb(s, l, t, w, h, lines, sz=16, c=INK, al=PP_ALIGN.LEFT):
    bx = s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h))
    tf = bx.text_frame; tf.word_wrap = True
    for i, li in enumerate(lines):
        p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
        if isinstance(li, dict):
            p.text = li.get('t',''); p.font.size = Pt(li.get('s',sz))
            p.font.color.rgb = li.get('c',c); p.font.bold = li.get('b',False)
            p.font.name = li.get('f','Inter'); p.alignment = li.get('a',al)
            p.space_before = Pt(li.get('sb',2)); p.space_after = Pt(li.get('sa',2))
        else:
            p.text = str(li); p.font.size = Pt(sz); p.font.color.rgb = c
            p.font.name = "Inter"; p.alignment = PP_ALIGN.LEFT; p.space_after = Pt(2)
    return bx

def bar(s, l=0.8, t=0.3, w=1.2, h=0.04):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(l),Inches(t),Inches(w),Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = ORANGE; sh.line.fill.background()
    return sh

def rect(s, l, t, w, h, c):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(l),Inches(t),Inches(w),Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = c; sh.line.fill.background()
    return sh

def pn():
    global N; N+=1; return N

def cover(txt, sub="", subtitle2="", quote_txt="", quote_by=""):
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, DARK)
    tb(s, 0.8, 1.5, 11.5, 1.8, txt, 72, WHITE, fn="Inter")
    if sub: tb(s, 0.8, 3.3, 11.5, 0.6, sub, 22, WHITE, fn="Noto Sans SC")
    bar(s, 0.8, 4.1, 2, 0.04)
    if subtitle2: tb(s, 0.8, 4.4, 11.5, 0.5, subtitle2, 16, G3, fn="Noto Sans SC")
    if quote_txt: tb(s, 0.8, 6.0, 11, 0.4, f'"{quote_txt}"', 13, G3, fn="Inter")
    tb(s, 11.5, 7.0, 1.5, 0.3, str(pn()), 10, G3, al=PP_ALIGN.RIGHT)

def content(t, items, q="", qa="", dark=False):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, DARK if dark else LIGHT)
    tc = WHITE if dark else INK
    bar(s, 0.8, 0.3, 1.2, 0.04)
    tb(s, 0.8, 0.5, 11.5, 0.6, t, 28, tc, fn="Noto Sans SC")
    rect(s, 0.8, 1.2, 11.5, 0.005, G2)
    mtb(s, 0.8, 1.6, 11.5, 4.2, items, 15, tc if not dark else WHITE)
    if q: tb(s, 0.8, 6.2, 11, 0.4, f'"{q}"', 12, G3, fn="Inter")
    tb(s, 11.5, 7.0, 1.5, 0.3, str(pn()), 10, G3, al=PP_ALIGN.RIGHT)

def data(t, pts, q="", qa=""):
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, DARK)
    bar(s, 0.8, 0.3, 1.2, 0.04)
    tb(s, 0.8, 0.5, 11.5, 0.6, t, 28, WHITE, fn="Noto Sans SC")
    rect(s, 0.8, 1.2, 11.5, 0.005, G2)
    cols = min(len(pts), 5); cw = 10.5 / cols
    for i, dp in enumerate(pts):
        x = 0.8 + (i%cols)*(cw+0.2); r = i//cols; y = 1.6 + r*2.3
        tb(s, x, y, cw, 0.9, dp.get('n',''), 52, ORANGE, b=True, fn="Inter")
        tb(s, x, y+0.9, cw, 0.35, dp.get('l',''), 14, G2, fn="Noto Sans SC")
        tb(s, x, y+1.25, cw, 0.3, dp.get('sl',''), 11, G3, fn="Inter")
    if q: tb(s, 0.8, 6.2, 11, 0.4, f'"{q}"', 12, G3, fn="Inter")
    tb(s, 11.5, 7.0, 1.5, 0.3, str(pn()), 10, G3, al=PP_ALIGN.RIGHT)

def tline(t, items, q="", qa=""):
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, DARK)
    bar(s, 0.8, 0.3, 1.2, 0.04)
    tb(s, 0.8, 0.5, 11.5, 0.6, t, 28, WHITE, fn="Noto Sans SC")
    rect(s, 0.8, 1.2, 11.5, 0.005, G2)
    y = 1.6
    for item in items:
        rect(s, 0.8, y+0.05, 0.12, 0.12, ORANGE)
        tb(s, 1.2, y-0.05, 2.0, 0.3, item.get('t',''), 13, ORANGE, b=True, fn="JetBrains Mono")
        tb(s, 3.5, y-0.05, 9, 0.4, item.get('d',''), 14, WHITE, fn="Inter")
        y += 0.55
    if q: tb(s, 0.8, 6.2, 11, 0.4, f'"{q}"', 12, G3, fn="Inter")
    tb(s, 11.5, 7.0, 1.5, 0.3, str(pn()), 10, G3, al=PP_ALIGN.RIGHT)

def act_title(n, title, duration, q="", qa=""):
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, DARK)
    tb(s, 0.8, 1.5, 3, 1.5, n, 96, ORANGE, b=True, fn="Inter")
    bar(s, 0.8, 3.2, 1.5, 0.04)
    tb(s, 0.8, 3.5, 11, 1.2, title, 36, WHITE, fn="Noto Sans SC")
    tb(s, 0.8, 5.0, 5, 0.4, duration, 14, G3, fn="Inter")
    if q: tb(s, 0.8, 6.0, 11, 0.4, f'"{q}"', 12, G3, fn="Inter")
    tb(s, 11.5, 7.0, 1.5, 0.3, str(pn()), 10, G3, al=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════════
# DECK CONTENT
# ══════════════════════════════════════════════════════════

# ── 01: Cover ──
cover("🧠 DereInside",
      "从本地知识库到认知引擎",
      "Know your project from the inside out.",
      "The best way to predict the future is to build it.", "Alan Kay")

# ── 02: What is DereInside ──
content("DereInside 是什么", [
    {"t":"一个本地优先的 AI 知识系统",'b':True,'s':18,'c':ORANGE},
    "  不是向量数据库的聊天包装器",
    "  是一个多模型认知引擎",
    "",
    {"t":"它不是检索工具 — 它理解你的代码",'b':True,'s':17},
    "",
    {"t":"三个核心特质",'b':True,'s':16,'c':ORANGE},
    "  🧬  多模型认知 — 6 个可互换的 AI 模型自动调度",
    "  🔗  知识图谱 — 实体/关系/跨域融合的活跃知识库",
    "  🧠  自学习共识 — 多模型交叉验证自动滤噪",
], "", "")

# ── 03: Why not a black box ──
data("为什么不是另一个黑盒", [
    {'n':'6+','l':'可互换模型','sl':'Ollama/vLLM/OpenAI/FreCode/MiniMax'},
    {'n':'4','l':'能力维度','sl':'智能×成本×速度×质量'},
    {'n':'5','l':'提取模式','sl':'从正则到7B-LLM, 自动选择'},
    {'n':'21','l':'知识域(Wings)','sl':'代码/文档/对话/运维...'},
], "Most AI knowledge tools are black boxes. One model. One pipeline. One size fits nobody.",
   "DereInside")

# ── 04: Act I Title ──
act_title("I","已经做到的：从架构到落地", "15 slides",
          "Architecture is not about building. It's about the decisions that make building possible.",
          "Unknown")

# ── 05: Model Registry ──
content("Model Registry — 一等公民的 AI 端点", [
    {"t":"每个 AI 端点都有四维画像：",'b':True,'s':16,'c':ORANGE},
    "  智能等级 (intelligence): low → medium → high → very_high",
    "  成本等级 (cost_tier): free → cheap → moderate → expensive",
    "  速度等级 (speed_tier): fast → medium → slow",
    "  质量等级 (quality): low → medium → high",
    "",
    {"t":"不是配置 — 是声明式能力声明",'b':True,'s':16,'c':ORANGE},
    "  Qwen-7b → Ollama, 免费, 高质, 慢速",
    "  GPT-4o → OpenAI, 昂贵, 极高质, 快速",
    "  FreeCode-Mini → FreeCode, 免费, 中质, 极速",
    "",
    {"t":"切换 provider 不改一行代码",'b':False,'s':15,'c':GREEN},
], "", "")

# ── 06: Pipeline Resolver ──
content("Pipeline Resolver — 约束求解，不是 fallback 链", [
    {"t":"传统做法：A 失败 → B 失败 → C",'b':True,'s':16,'c':G3},
    "",
    {"t":"DereInside 的做法：",'b':True,'s':16,'c':ORANGE},
    "  1️⃣ 过滤：智能等级 ≥ X, 成本 ≤ Y, 延迟 ≤ Z",
    "  2️⃣ 健康检查：模型在线 + 响应正常",
    "  3️⃣ 排序：按目标优化（质量优先 / 速度优先 / 成本最低）",
    "  4️⃣ 降级：无匹配项时放松约束，用日志告知",
    "",
    {"t":"例子：代码实体提取",'b':True,'s':15,'c':BLUE},
    "  extract 管道要求: {min_intelligence: low, max_cost: free, max_latency_ms: 10000}",
    "  候选: Qwen-7b ✓ / GPT-4o ✗ (太贵) / FreeCode ✗ (不在候选列表)",
], "", "")

# ── 07: Model Profiler ──
content("Model Profiler — 零配置自动探测", [
    {"t":"用户不需要知道模型的规格参数",'b':True,'s':16,'c':ORANGE},
    "",
    "  Boot Probe: 3 个 golden sample, ~2s 定位模型能力基线",
    "  Deep Probe: 15 个样本, ~10s 精确评分 (智能/速度/质量)",
    "  Cost 推定: 基于驱动器和 URL 自动判断 (Ollama=free, OpenAI=expensive)",
    "",
    {"t":"被动观察者 (Passive Observer)",'b':True,'s':16,'c':ORANGE},
    "  运行时零成本收集: 延迟 / 实体数量 / 错误率",
    "  检测异常: 延迟翻倍 / 实体数骤降 → 自动触发重新 profiling",
    "",
    {"t":"震荡检测",'b':False,'s':15,'c':RED},
    "  FreeCode 模型反复跳变 → 3 次探测后冻结配置并告警",
], "", "")

# ── 08: 5-Mode Extraction ──
content("5 种提取模式 — 自动调度", [
    {"t":"每个 chunk 自动匹配最优提取模式：",'b':True,'s':16,'c':ORANGE},
    "",
    "  .java/.py → regex (精度 95.6%, 5ms)",
    "  .md/.txt  → hybrid-1.5b (F1=0.43, 3s)",
    "  .xml/.sql → 1.5b (召回 55.9%, 2.6s)",
    "  .log      → skip (零成本)",
    "  通用     → hybrid-7b (平衡模式)",
    "",
    {"t":"效果",'b':True,'s':16,'c':GREEN},
    "  2,467 chunks 从 5h26m 降低到 ~1.5h",
    "  加权 F1 提升约 15%",
], "", "")

# ── 09: Self-Learning Consensus ──
content("Self-Learning Consensus — 自学习共识", [
    {"t":"不信任单一模型 — 多模式交叉验证",'b':True,'s':16,'c':ORANGE},
    "",
    "  Chunk: \"class OrderService extends BaseService { @Autowired ... }\"",
    "",
    "  regex     → {OrderService, BaseService}            (精 95.6%)",
    "  hybrid-7b → {OrderService, BaseService, ...}        (精 77.4%)",
    "  1.5b      → {OrderService, BaseService, Autowired}  (高召回, 有噪音)",
    "",
    "  ConsensusEngine:",
    "    OrderService  3/3 → confirmed (weight=1.0) ✓",
    "    BaseService   3/3 → confirmed (weight=1.0) ✓",
    "    Autowired     1/3 → rejected  (weight=0.0) ✗ (噪音: 非实体)",
    "",
    {"t":"结果：噪音从 ~28% 降至 ~10%（3 次完整循环后）",'b':True,'s':15,'c':GREEN},
], "", "")

# ── 10: Knowledge Graph ──
data("知识图谱 — 5,485 实体 · 10,817 链接", [
    {'n':'5,485','l':'知识图谱实体','sl':'类/接口/概念/文档'},
    {'n':'10,817','l':'关系链接','sl':'extends/depends_on/serves_path'},
    {'n':'7.9x','l':'实体提升','sl':'对比旧单模型架构'},
    {'n':'~15%','l':'F1 提升','sl':'加权后对比硬编码'},
], "The knowledge graph explosion is the direct result of multi-model architecture — not a better single model.",
   "aITMS01")

# ── 11: Entity Resolution ──
content("实体消歧与跨域融合", [
    {"t":"解决最棘手的问题：同一实体不同命名",'b':True,'s':16,'c':ORANGE},
    "",
    "  别名字典: KYCApplication = kycapplication = KYC 申请 → 自动合并",
    "  后缀剥离: AuditServiceImpl → AuditService",
    "  跨 Wing 融合: 不同知识域的同名实体自动链接",
    "",
    {"t":"子图查询 — 交互式知识探索",'b':True,'s':16,'c':ORANGE},
    "  derekinside graph subgraph \"KYC\" --depth 2 --ascii",
    "  → KYC (concept)",
    "      ├─ KYCApplication (class) ← depends_on",
    "      ├─ KYCController (class) ← serves_path: /api/kyc/submit",
    "      └─ 合同审批流程 (concept) ← related",
], "", "")

# ── 12: Smart Dispatch ──
content("Smart Dispatch — 智能调度", [
    {"t":"不是所有 content 都需要大模型",'b':True,'s':16,'c':ORANGE},
    "",
    "  代码文件 → regex (95.6% 精度, 5ms) — 完全不需要 LLM",
    "  自然语言文档 → hybrid-1.5b (F1=0.43, 3s) — 轻量模型就够",
    "  复杂配置文件 → 1.5b (召回 55.9%, 2.6s) — 高召回优先",
    "  日志文件 → skip (0s) — 不是知识",
    "",
    {"t":"分类器自动判断内容类型",'b':True,'s':15,'c':BLUE},
    "  基于文件名 + 首行 + 结构特征 → 标签分类",
    "  准确率: ~94% (基于 599 pages / 2,931 chunks 生产数据)",
], "", "")

# ── 13: Quantified Benchmarks ──
data("LongMemEval — 量化基准测试", [
    {'n':'95.6%','l':'regex 精度','sl':'290-entity golden dataset'},
    {'n':'59.0%','l':'7b LLM 精度','sl':'全模型模式'},
    {'n':'72.4%','l':'hybrid-1.5b 精度','sl':'轻量混模式'},
    {'n':'55.9%','l':'1.5b 召回','sl':'最高召回模式'},
], "Every mode benchmarked. No black boxes.", "LongMemEval")

# ── 14: MCP Integration ──
content("MCP 集成 — Agent 原生认知", [
    {"t":"不只是 REST API — 是 Agent 的认知管道",'b':True,'s':16,'c':ORANGE},
    "",
    "  JSON-RPC 2.0 over stdio — 零网络延迟",
    "  FastAPI HTTP bridge (port 18890) — 远程访问",
    "  Per-Agent 隔离命名空间 — 每个 Agent 有独立搜索上下文",
    "",
    {"t":"子 Agent 启动时自动注入",'b':True,'s':16,'c':ORANGE},
    "  开工前: derekinside search → 获取相关实体 + 关系 + 排序的 chunk",
    "  开工后: 每条决策日志 → 后续 Agent 可检索",
    "",
    {"t":"搜索响应: <200ms 典型延迟",'b':True,'s':15,'c':GREEN},
], "", "")

# ── 15: Act I — Production Results ──
data("生产运行数据 — 零云依赖", [
    {'n':'21','l':'知识域(Wings)','sl':'对比初始的14'},
    {'n':'2,931','l':'索引块(Chunks)','sl':'100% 已嵌入'},
    {'n':'7.9x','l':'实体增长','sl':'旧模型架构vs新架构'},
    {'n':'1.5h','l':'全索引时间','sl':'之前5h26m'},
], "One VM. PostgreSQL. Zero cloud API calls.", "Production")

# ── 16: Act II Title ──
act_title("II","准备实现的：Roadmap", "10 slides",
          "Vision without execution is hallucination.", "Thomas Edison")

# ── 17: Phase 4 — Multi-Model Ensemble ──
content("Phase 4: 多模型集成 + Agent 上下文门", [
    {"t":"Multi-Model Ensemble",'b':True,'s':16,'c':ORANGE},
    "  不是选择一个模型，而是让多个模型投票",
    "  提取时: 2-3 个模型同时运行 → 加权投票 → 最佳答案",
    "  搜索时: 多模型 rerank → 消除单一模型偏差",
    "  代价: 延迟增加 2-3x, 精度提升预估 8-12%",
    "",
    {"t":"Agent-Native Context Gate",'b':True,'s':16,'c':ORANGE},
    "  子 Agent 启动时自动组装结构化的知识上下文",
    "  不再只有 \"这里有一些相关文档\"",
    "  而是: \"实体A的3个关系路径 + 最近决策日志 + 活跃事实\"",
], "", "")

# ── 18: Phase 5 — Web UI ──
content("Phase 5: Web UI Dashboard", [
    {"t":"不是管理面板 — 是知识指挥中心",'b':True,'s':16,'c':ORANGE},
    "",
    "  搜索界面: 带图谱预览 + 关联推荐 + 语义过滤",
    "  知识图谱可视化: 力导向图 / 层次树 / 子图聚焦",
    "  状态面板: Wing 热度 / 嵌入进度 / 提取质量 / 模型健康",
    "  协作注释: 允许用户在实体上添加 + 分享标签 / 说明",
    "",
    {"t":"技术方案",'b':True,'s':15,'c':BLUE},
    "  Vue 3 + D3.js (力导向图) 或 Cytoscape.js (层次树)",
    "  REST API 已就绪 — 纯前端增量开发",
], "", "")

# ── 19: Phase 6 — Fleet Learning ──
content("Phase 6: Fleet Learning", [
    {"t":"一个 profle 被共享 = 所有实例受益",'b':True,'s':16,'c':ORANGE},
    "",
    "  场景: 你在本地跑了 5 次 profile, 发现了 Qwen-7b 的最优参数",
    "  签入 git: 自动合并到 Fleet 配置树",
    "  其他实例下次 pull: 拿到你的 profile 结果, 跳过冷启动",
    "",
    {"t":"好处",'b':True,'s':15,'c':GREEN},
    "  分布式模型评测 → 社区基准数据库",
    "  \"这个模型在 Python 代码提取上不如 regex\" — 数据说话",
    "  新模型接入后 3 天内获得社区 profilng 结果",
], "", "")

# ── 20: RFC-0001 — Fact Logging ──
content("RFC-0001: Agent 事实日志与共享记忆", [
    {"t":"问题：子 Agent 之间没有结构化共享内存",'b':True,'s':16,'c':ORANGE},
    "  每次 sessions_spawn 看到的都是干净 session",
    "  前一个 Agent 的知识丢失 → 重复劳动 / 矛盾决策",
    "",
    {"t":"设计",'b':True,'s':16,'c':BLUE},
    "  Fact 数据模型: {id, fact_text, valid_from/to, agent_id, episode_type}",
    "  新增 facts + fact_entities 表 (PostgreSQL)",
    "  POST /api/v1/facts — 记录事实",
    "  POST /api/v1/facts/search — 时序搜索",
    "",
    {"t":"Agent 指令模式",'b':True,'s':15,'c':GREEN},
    "  search before ask: 开工前先查事实",
    "  log after decision: 关键决策后记事实",
    "", {"t":"工作量: ~235 行 · 低-中难度",'b':True,'s':14,'c':ORANGE},
], "", "")

# ── 21: RFC-0002 — EverOS Merger ──
data("RFC-0002: DereInside × EverOS 合并方案", [
    {'n':'~8k','l':'DereInside 代码','sl':'Python + PostgreSQL'},
    {'n':'~33k','l':'EverOS 代码','sl':'Python + SQLite + LanceDB'},
    {'n':'1-2d','l':'路径A: 依赖式合并','sl':'0 代码重写, 复用33k行'},
    {'n':'~18d','l':'路径B: 吸收式合并','sl':'白己实现核心功能'},
], "建议 Phase 1 → 路径A (1-2天拿90%价值), Phase 2 → 逐步按需吸收子模块",
   "aITMS01")

# ── 22: Roadmap Summary ──
data("Roadmap 全景", [
    {'n':'4','l':'Phase 4 — 🚧','sl':'多模型集成 + 上下文门'},
    {'n':'5','l':'Phase 5 — 📋','sl':'Web UI Dashboard'},
    {'n':'6','l':'Phase 6 — 📋','sl':'Fleet Learning'},
    {'n':'RFC','l':'事实日志 + EverOS','sl':'见 RFC-0001/0002'},
], "有一就有二。所有路线图节点都在 GitHub RFC 中 — 不是墙上贴纸。", "aITMS01")

# ── 23: Summary ──
cover("架构问答",
      "\"One embedding model for everything\" is a solved problem.",
      "DereInside asks better questions.",
      "The important thing is not to stop questioning.", "Albert Einstein")

# ── 24: Docs ──
content("文档与资源", [
    {"t":"GitHub: github.com/derekwang85/derekinside",'b':True,'s':18,'c':BLUE},
    "",
    "  README: 架构概览 + Quick Start + Benchmarks",
    "  RFC-0001: Agent Fact Logging & Shared Memory",
    "  RFC-0002: DereInside × EverOS Merge Analysis",
    "  docs/: 完整接口文档 + 配置指南 + 部署手册",
    "",
    {"t":"运行状态",'b':True,'s':16,'c':GREEN},
    "  进程: derekinside serve --mode http --host 0.0.0.0 --port 18890",
    "  状态: βeta · 可用 · 零云依赖",
    "  API: GET /health, GET /api/v1/status, POST /api/v1/search",
], "", "")

# ── 25: Q&A ──
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, ORANGE)
tb(s, 0.8, 2.0, 11, 1.5, "Q & A", 72, WHITE, fn="Inter")
rect(s, 0.8, 3.8, 2, 0.04, WHITE)
tb(s, 0.8, 4.2, 11, 1.0, "Derek · aITMS01 — 2026", 20, WHITE, fn="Noto Sans SC")
tb(s, 11.5, 7.0, 1.5, 0.3, str(pn()), 10, WHITE, fn="JetBrains Mono", al=PP_ALIGN.RIGHT)

# ── 26: Thanks ──
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, DARK)
tb(s, 0.8, 2.5, 11, 1.5, "Thanks", 64, WHITE, fn="Inter")
tb(s, 0.8, 4.2, 11, 0.5, "github.com/derekwang85/derekinside", 16, G3, fn="JetBrains Mono")
tb(s, 11.5, 7.0, 1.5, 0.3, str(pn()), 10, G3, fn="JetBrains Mono", al=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════════════════════
out = os.path.expanduser("~/derekinside/DereInside-2026.pptx")
prs.save(out)
print(f"\n✅ Saved {out} ({N} slides)")
